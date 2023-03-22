from pptx import Presentation


def pptx_to_text(input_filename: str, slides_seperator: str = "\n\n") -> str:
    presentation = Presentation(input_filename)
    presentation_text = ""

    for slide in presentation.slides:
        
        slide_has_title = slide.shapes.title is not None

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue

            shape_text = f'\n{shape.text}'

            if slide_has_title and shape.text == slide.shapes.title.text:
                shape_text += ":"

            presentation_text += shape_text

        presentation_text += slides_seperator

    return presentation_text
